from pptx import Presentation

def extract_text_from_pptx(file_path):
    ppt = Presentation(file_path)
    text_runs = []

    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame
                if text_frame:
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
    return text_runs

# Example usage:
file_path = 'D:/Users/Christan/Downloads/ENG М1L1. Introduction  to Unity.pptx'
all_text = extract_text_from_pptx(file_path)
print(all_text)
