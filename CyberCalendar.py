import re
import datetime
from pptx import Presentation
from pptx.util import Inches

def read_entries(filename):
    """
    Reads the text file and splits it into entries based on one or more blank lines.
    Removes any lines that are just numbers.
    """
    with open(filename, 'r', encoding='utf-8') as f:
        text = f.read()

    # Split text into entries separated by one or more blank lines
    entries = [entry.strip() for entry in re.split(r'\n\s*\n', text) if entry.strip()]

    processed_entries = []
    for entry in entries:
        # Remove lines that are just numbers
        lines = entry.split('\n')
        lines = [line for line in lines if not line.strip().isdigit()]
        processed_entry = '\n'.join(lines).strip()
        if processed_entry:
            processed_entries.append(processed_entry)
    return processed_entries

def assign_dates(entries, start_date):
    """
    Assigns dates to each entry starting from the given start_date.
    """
    dates = [start_date + datetime.timedelta(days=i) for i in range(len(entries))]
    return dates

def create_presentation(entries, dates, output_file):
    """
    Creates a PowerPoint presentation with one slide per entry.
    Each slide contains the date and the corresponding text.
    """
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Choosing a blank slide layout

    for entry, date in zip(entries, dates):
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add date as the title
        title_placeholder = slide.shapes.title
        if not title_placeholder:
            title_placeholder = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_placeholder.text = date.strftime('%B %d, %Y')

        # Add entry text
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = entry

    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

def main():
    input_filename = 'Cyberkalender.txt'
    output_filename = 'Cyberkalender.pptx'

    entries = read_entries(input_filename)
    start_date = datetime.date(2025, 1, 1)
    dates = assign_dates(entries, start_date)
    create_presentation(entries, dates, output_filename)

if __name__ == '__main__':
    main()
